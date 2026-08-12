"""The deck, in three sections.

It opens with a run-parameters slide - when the simulations were produced and
every input they were produced from - so a deck found on disk months later can
still be read against the assumptions behind it.

**1. Optimal stake distribution** - which stakes are worth playing and in what
proportion: the stake screen, the bb and euro waterfalls, the frontier, the
optimum with its nearest alternatives, and a simulation of it.

**2. Shot-taking** - the two ways to seat a table on the rung above the optimal
mix (paying for it from the top of the mix, or from the bottom) and what each
costs in risk.

**3. Current configuration** - the mix actually played, reconstructed from hands
per stake and simulated the same way as the optimum. It has no table slide of
its own: it is the first row of EVERY table in the deck, highlighted orange
against the optimum's green, so each table answers 'what does this change,
against what I do today' without the reader carrying a number between slides.

Then two appendices: why a mix can earn more and risk less, and the method.

Each section's simulation sits in the section whose mix it simulates, so the
optimum and the real allocation are never compared across a section break.

Every number is recomputed here from the same functions the CLI prints, so a
slide cannot disagree with the terminal. Nothing is hardcoded.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from . import charts, estimation, mix, pptx_common as pc, progress, rates, sim  # noqa: E402
from .config import Config  # noqa: E402

__all__ = ["build"]

# Waterfall colours: money in, money out, and the two totals.
COL_POSITIVE = "#1baf7a"
COL_NEGATIVE = "#d03b3b"
COL_TOTAL = "#2a78d6"
COL_OPENING = "#b6b5b1"
"""The opening bar. Grey rather than the closing bar's blue: it is the starting
position, not a result, and colouring both ends alike invited them to be read as
the same kind of quantity. Blue is now reserved for the number you keep."""
COL_INK = "#0b0b0b"
COL_MUTED = "#7a7972"


def _money(config: Config, eur: float, dp: int = 0) -> str:
    """A euro amount, labelled in the display currency: `GBP 8,621`.

    Every money string on every slide goes through here or `_plain`, so a deck
    built in sterling has no euro figures hiding in a footnote. Takes EUROS - the
    internal unit - so nothing can be converted twice.
    """
    return config.currency.fmt(eur, dp)


def _plain(config: Config, eur: float, dp: int = 0) -> str:
    """As `_money`, without the code - for table cells whose header carries it."""
    return config.currency.plain(eur, dp)


def risk_rule_note(config: Config) -> str:
    """One sentence naming the rule that actually chose the mixes on this deck."""
    from . import tolerance as _tolerance

    described = _tolerance.for_config(config).describe(config)
    if config.risk_mode == "both":
        return (
            f"THE RULE APPLIED: {described}. Two constraints, and the chosen mix had to clear "
            f"BOTH - so whichever is stricter at this bankroll is the one that decided. They "
            f"are not two spellings of one idea: ruin asks whether you survive, while the "
            f"downswing bar asks how deep a PEAK-TO-TROUGH fall you sit through on the way, "
            f"measured from whatever high the roll had reached rather than from where it "
            f"started. A mix can be certain to survive and still be unplayable."
        )
    if config.risk_mode == "downswing":
        return (
            f"THE RULE APPLIED: {described}. That is a "
            f"PEAK-TO-TROUGH fall - measured from whatever high the bankroll had reached, not "
            f"from where it started - and it is the constraint the chosen mix had to satisfy. "
            f"Risk of ruin is still reported throughout, but it did not decide anything here."
        )
    return (
        f"THE RULE APPLIED: {described}. Downswing columns "
        f"are reported for context; they did not gate the choice."
    )


def drawdown_note(config: Config) -> str:
    """Defined once and repeated under every table carrying the columns - a
    reader landing on one slide should not have to hunt for what they mean."""
    hands = config.timescale_hands
    hours = hands / (config.tables * config.hands_per_hour_per_table)
    paths = sim.TABLE_PATHS
    return (
        f"HOW THE DOWNSWING COLUMNS ARE CALCULATED.  (1) Simulate one lifetime: {hands:,} "
        f"hands ({hours:,.0f} hours at {config.tables} tables) played at that exact mix. "
        f"(2) Within that lifetime find the single deepest PEAK-TO-TROUGH fall - how far the "
        f"bankroll dropped below whatever high it had previously reached. That is one number "
        f"per lifetime. (3) Repeat for {paths:,} independent lifetimes. (4) Report three "
        f"points on the resulting distribution: MEDIAN is the middle one, so half of "
        f"lifetimes fall deeper than that and half do not; 10% WORST is the 90th percentile, "
        f"one lifetime in ten; 1% WORST is the 99th, one in a hundred.  Note these are NOT "
        f"percentiles of every downswing you have - most of those are trivial and there are "
        f"thousands of them. Each lifetime contributes only its WORST one. The timescale is "
        f"part of the number: given unlimited time a peak-to-trough fall grows without bound, "
        f"because a winning bankroll keeps making new highs to fall from. Change the "
        f"timescale and these change with it.  ON TABLES is a different kind of figure - not "
        f"a rate or a probability but a stock: the money sitting in front of you at once, at "
        f"a full 100bb buy-in per seat. It is what a mix has at stake today rather than "
        f"eventually, and the number to hold against the bankroll behind it."
    )


def _fit_columns(columns):
    """Scale a column spec down if it would run off the right-hand edge.

    A table wider than the content area is centred by the callers, which puts
    half the overflow off each side and silently truncates BOTH ends - the
    failure is invisible in code review and obvious only on the rendered slide.
    Scaling proportionally keeps the relative widths the spec asked for and
    guarantees the table fits, whatever gets added to it later.
    """
    limit = pc.CONTENT_WIDTH / Inches(1)
    total = sum(width for _, width in columns)
    if total <= limit:
        return columns
    factor = limit / total
    return [(label, width * factor) for label, width in columns]


def timescale_label(config: Config) -> str:
    """The simulation horizon as a short header word - '1M', '200k'."""
    hands = config.timescale_hands
    if hands >= 1_000_000 and hands % 1_000_000 == 0:
        return f"{hands // 1_000_000}M"
    if hands >= 1_000 and hands % 1_000 == 0:
        return f"{hands // 1_000}k"
    return f"{hands:,}"


CURRENT_LABEL = "CURRENT (July)"
"""What the played mix is called wherever it appears. One constant, because it
now appears on four tables and a caption that disagreed with the others would
read as a different mix."""


def _benchmark_row(current):
    """The played mix, as the first row of an allocation table.

    Every table opens with it so each one answers the same question - what does
    this change, against what is actually being done today - without the reader
    holding a number in their head from an earlier slide."""
    return (current, "what you play now", pc.COL_ORANGE, True)


def horizon_ev(config: Config, allocation) -> float:
    """Expected profit over the deck's horizon, in euros.

    Linear in hands - nothing compounds, since the mix is static - so this is
    exactly the EUR/hr column on a different clock. It earns its place because
    a year's money is the figure people actually weigh; an hourly rate is not.
    """
    return allocation.mean_eur_per_100 * config.timescale_hands / 100.0


# --------------------------------------------------------------------------- #
# Cover - when this was run, and on what
# --------------------------------------------------------------------------- #
def _run_parameters_slide(prs, layouts, config: Config):
    """Every input the deck was produced from, on one slide.

    Read off the live Config, which is the object the simulations actually ran
    on - so a CLI override like `--bankroll 20000` shows the value that was
    used, not the value sitting in config.toml. A deck whose parameters slide
    could disagree with its own numbers would be worse than no slide."""
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "How this deck was produced")

    built = datetime.now().astimezone()
    stamp = slide.shapes.add_textbox(
        pc.CONTENT_LEFT, pc.CONTENT_TOP - Inches(0.05), pc.CONTENT_WIDTH, Inches(0.35)
    )
    run = stamp.text_frame.paragraphs[0].add_run()
    run.text = f"Simulations run {built:%A %d %B %Y at %H:%M %Z}"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = pc.BLACK

    total_hands_per_hour = config.tables * config.hands_per_hour_per_table
    timescale_hours = config.timescale_hands / total_hands_per_hour
    downswing = (
        f"{config.downswing_probability:.0%} chance of "
        f"{_money(config, config.downswing_amount_eur)} in "
        f"{config.downswing_hands:,} hands"
        if config.downswing_amount_eur is not None
        else "not set"
    )
    settings = [
        ("Bankroll", _money(config, config.bankroll_eur)),
        ("Tables played at once", f"{config.tables}"),
        ("Risk rule applied", f"{config.risk_mode}"),
        ("Risk-of-ruin tolerance", f"{config.ruin_tolerance:.2%}"),
        ("Downswing tolerance", downswing),
        ("Rakeback", f"{config.rakeback_pct:.0%} of rake paid"),
        ("Hands/hour/table", f"{config.hands_per_hour_per_table:,.0f}"
                             f"  ({total_hands_per_hour:,.0f} total)"),
        ("Simulation timescale", f"{config.timescale_hands:,} hands"
                                 f"  (~{timescale_hours:,.0f} hrs)"),
        ("Lifetimes simulated", f"{config.sim_paths:,}"),
        ("Kelly fraction", f"{config.kelly_fraction:g}  (report/stake only)"),
        ("Win-rate haircut per table", f"{config.winrate_haircut_bb_per_table:.2f} bb/100"),
        ("Table correlation", f"{config.table_correlation:.2f}"),
        ("Money shown in",
         config.currency.code if config.currency.is_base
         else f"{config.currency.code}  (fixed "
              f"{config.currency.eur_per_unit:g} EUR = 1 {config.currency.code})"),
    ]

    top = pc.CONTENT_TOP + Inches(0.42)
    left_cols = [("Setting", 3.00), ("Value", 2.20)]
    left_width = Inches(sum(w for _, w in left_cols))
    shape = slide.shapes.add_table(
        len(settings) + 1, len(left_cols), pc.CONTENT_LEFT, top,
        left_width, Inches(0.34 + 0.30 * len(settings)),
    )
    table = shape.table
    for index, (_, width) in enumerate(left_cols):
        table.columns[index].width = Inches(width)
    for index, (label, _) in enumerate(left_cols):
        cell = table.cell(0, index)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=9, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE)
    for row, (label, value) in enumerate(settings, start=1):
        for index, text in enumerate((label, value)):
            cell = table.cell(row, index)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, text, font_size=9, bold=(index == 0), font_colour=pc.BLACK)

    # ---- the per-stake inputs, every column config.toml can carry ---------- #
    stake_cols = [
        ("Stake", 0.75),
        ("bb\nEUR", 0.70),
        ("Win rate\nbb/100", 0.95),
        ("SD\nbb/100", 0.80),
        ("Rake\nbb/100", 0.85),
        ("Sample\nhands", 1.00),
        ("Hands\nplayed", 1.00),
        ("Max\ntables", 0.85),
    ]
    stake_width = Inches(sum(w for _, w in stake_cols))
    stake_left = int(pc.CONTENT_LEFT + pc.CONTENT_WIDTH - stake_width)
    shape = slide.shapes.add_table(
        len(config.stakes) + 1, len(stake_cols), stake_left, top,
        stake_width, Inches(0.40 + 0.30 * len(config.stakes)),
    )
    table = shape.table
    for index, (_, width) in enumerate(stake_cols):
        table.columns[index].width = Inches(width)
    for index, (label, _) in enumerate(stake_cols):
        cell = table.cell(0, index)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=9, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE, wrap=True)

    for row, stake in enumerate(config.stakes, start=1):
        # An omitted optional key is shown as what it MEANS, not as a blank: a
        # missing max_tables is "no limit", and a missing sample is "-" because
        # the win rate is then a guess with no interval behind it.
        values = [
            stake.name,
            f"{stake.bb_eur:.2f}",
            f"{stake.winrate_bb100:.2f}",
            f"{stake.stdev_bb100:.1f}",
            f"{stake.rake_bb100:.3f}" if stake.rake_bb100 is not None else "-",
            f"{stake.hands:,}" if stake.hands else "-",
            f"{stake.current_hands:,.0f}" if stake.current_hands is not None else "-",
            f"{stake.max_tables}" if stake.max_tables is not None else "no limit",
        ]
        for index, value in enumerate(values):
            cell = table.cell(row, index)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, value, font_size=9, bold=(index == 0), font_colour=pc.BLACK)

    _footnote(
        slide,
        "These are the values the run actually used, including any command-line override, "
        "so they cannot disagree with the numbers on the slides that follow. 'Sample hands' "
        "is the volume the win rate is measured over and only drives the confidence interval. "
        "'Hands played' is the recent split used to reconstruct the current configuration, and "
        "is read as a ratio only. Standard deviation is an assumption at every stake, not a "
        "measurement.",
    )
    return slide


# --------------------------------------------------------------------------- #
# Slide 1 - the stake table
# --------------------------------------------------------------------------- #
def _stake_table_slide(prs, layouts, config: Config, screens: list[mix.StakeScreen]):
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "Every stake, priced on its own")

    columns = [
        ("Stake", 0.85),
        ("Hands", 0.95),
        ("Win rate\nbb/100", 0.95),
        ("Rake\nbb/100", 0.85),
        ("Rakeback\nbb/100", 0.95),
        ("Banked\nbb/100", 0.95),
        ("+/- 95%\nbb/100", 0.90),
        (f"Banked\n{config.currency.code}/100", 0.95),
        (f"SD\n{config.currency.code}/100", 0.90),
        (f"{config.currency.code}/hr", 0.85),
        (f"On tables\n{config.currency.code}", 0.95),
        ("Verdict", 2.55),
    ]
    rows = len(screens) + 1
    table_width = Inches(sum(w for _, w in columns))
    left = int(pc.CONTENT_LEFT + (pc.CONTENT_WIDTH - table_width) / 2)
    shape = slide.shapes.add_table(
        rows, len(columns), left, pc.CONTENT_TOP + Inches(0.35),
        table_width, Inches(0.4 + 0.32 * len(screens)),
    )
    table = shape.table
    for index, (_, width) in enumerate(columns):
        table.columns[index].width = Inches(width)

    for index, (label, _) in enumerate(columns):
        cell = table.cell(0, index)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=9, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE, wrap=True)

    for row, screen in enumerate(screens, start=1):
        stake = screen.stake
        rakeback = rates.rakeback_bb100(stake.rake_bb100, config.rakeback_pct)
        banked_bb = screen.mean_eur_per_100 / stake.bb_eur
        margin = (
            f"{estimation.winrate_ci(0.0, stake.stdev_bb100, stake.hands)[1]:.1f}"
            if stake.hands else "-"
        )
        # Excluded rows are greyed wholesale rather than flagged in one column:
        # the point is that none of their numbers feed the answer.
        excluded = not screen.kept
        ink = pc.GRID_GREY if excluded else pc.BLACK
        values = [
            stake.name,
            f"{stake.hands:,}" if stake.hands else "-",
            f"{stake.winrate_bb100:.2f}",
            f"-{stake.rake_bb100:.2f}" if stake.rake_bb100 else "-",
            f"+{rakeback:.2f}" if rakeback else "-",
            f"{banked_bb:.2f}",
            f"+/-{margin}",
            _plain(config, screen.mean_eur_per_100, 2),
            _plain(config, screen.stdev_eur_per_100),
            _plain(config, screen.eur_per_hour),
            _plain(config, screen.exposure_eur),
            "in the mix" if screen.kept else screen.excluded_reason,
        ]
        for index, value in enumerate(values):
            cell = table.cell(row, index)
            pc._zero_cell_margins(cell)
            pc.set_cell(
                cell, value, font_size=9, bold=(index == 0),
                bg_colour=pc.TABLE_LABEL_BG if excluded else None,
                # The verdict column, wherever it has ended up - an index literal
                # is what breaks silently when a column is inserted before it.
                font_colour=(
                    pc.COL_TEXT_RED if (excluded and index == len(values) - 1) else ink
                ),
            )

    _footnote(
        slide,
        "Win rate is net of rake, as the tracker reports it; rake and rakeback are shown "
        "separately so the composition is visible. Rakeback at "
        f"{config.rakeback_pct:.0%}. The +/- column is the 95% interval half-width implied "
        "by the sample - a win rate with no volume behind it is a guess, not a measurement.",
    )
    return slide


def _winrate_ci_slide(prs, layouts, config: Config, screens):
    """The edge and how well it is known, in both units, side by side."""
    slide = pc.add_image_slide(
        prs, layouts, "How well do you actually know each win rate?",
        charts.winrate_ci_figure(screens, config),
    )
    # add_image_slide centres the picture in the content area, which on a chart
    # this tall runs it into the footnote. Pin it to the top instead.
    picture = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    picture.top = pc.CONTENT_TOP

    assumed = [s.stake for s in screens if s.stake.measured_winrate_bb100 is not None]
    gap = ""
    if assumed:
        worst = max(
            assumed,
            key=lambda s: estimation.winrate_stderr(s.stdev_bb100, s.hands) if s.hands else 0.0,
        )
        half = estimation.Z_95 * estimation.winrate_stderr(worst.stdev_bb100, worst.hands)
        gap = (
            f" At {', '.join(s.name for s in assumed)} the model uses an ASSUMED rate, drawn as "
            f"the bar, while the diamond is what the sample says - {worst.name} measures "
            f"{worst.measured_winrate_bb100:.2f} bb/100 over {worst.hands:,} hands, "
            f"+/-{half:.1f} at 95%. The interval belongs to the diamond, never to the bar."
        )
    _footnote(
        slide,
        "Win rates are ALL-IN ADJUSTED - the luck is stripped from the spots where the money "
        "went in before the cards ran out, so the estimate converges faster than won-at-"
        "showdown. Rakeback carries no interval: it is a rebate on volume, known in advance."
        + gap,
    )
    return slide


# --------------------------------------------------------------------------- #
# Slides 2 and 3 - the waterfalls
# --------------------------------------------------------------------------- #
def _waterfall_figure(screens, config: Config, in_euros: bool):
    """One waterfall panel per kept stake, sharing a y axis.

    Bars: win rate before rake -> rake paid -> rakeback -> banked. Sharing the
    axis is the whole point of the slide; per-panel scaling would hide that the
    bars shrink in bb terms and grow in euros as the stakes rise.
    """
    kept = [s for s in screens if s.kept]
    charts._style()

    series = []
    for screen in kept:
        stake = screen.stake
        rake = stake.rake_bb100 or 0.0
        rakeback = rates.rakeback_bb100(stake.rake_bb100, config.rakeback_pct)
        before = stake.winrate_bb100 + rake
        banked = screen.mean_eur_per_100 / stake.bb_eur
        # bb -> money per 100 hands, with the display conversion folded in, so
        # the panel is drawn in the same units its axis label claims.
        scale = config.currency.from_eur(stake.bb_eur) if in_euros else 1.0
        series.append(
            (stake.name, before * scale, -rake * scale, rakeback * scale, banked * scale)
        )

    # The tallest thing on any panel is the opening or closing bar; the rake bar
    # hangs DOWN from the opening one, so it never sets the ceiling. (Its stored
    # height is already negative - adding it here once cost a third of the axis
    # to dead space.)
    y_max = max(max(before, banked) for _, before, _, _, banked in series) * 1.16
    y_min = min(0.0, min(banked for *_, banked in series)) * 1.1

    fig, axes = plt.subplots(
        1, len(series), figsize=(3.05 * len(series), 4.6), sharey=True, squeeze=False
    )
    unit = config.currency.axis("/ 100 hands") if in_euros else "bb / 100 hands"

    for ax, (name, before, rake, rakeback, banked) in zip(axes[0], series):
        after_rake = before + rake
        steps = [
            ("Before\nrake", 0.0, before, COL_OPENING),
            ("Rake", after_rake, -rake, COL_NEGATIVE),
            ("Rakeback", after_rake, rakeback, COL_POSITIVE),
            ("Banked", 0.0, banked, COL_TOTAL),
        ]
        for position, (label, bottom, height, colour) in enumerate(steps):
            ax.bar(position, height, bottom=bottom, color=colour, width=0.62,
                   edgecolor=charts.SURFACE, linewidth=1.5)
            value = before if position == 0 else (rake if position == 1 else
                                                  rakeback if position == 2 else banked)
            sign = "" if position in (0, 3) else ("-" if position == 1 else "+")
            ax.annotate(
                f"{sign}{abs(value):,.2f}" if not in_euros else f"{sign}{abs(value):,.1f}",
                xy=(position, bottom + height),
                xytext=(0, 4), textcoords="offset points",
                ha="center", fontsize=8.5, color=COL_INK, fontweight="bold",
            )
        # Connectors, so the eye follows the running total across the bars.
        ax.plot([0.31, 0.69], [before, before], color=COL_MUTED, linewidth=0.9, linestyle=":")
        ax.plot([1.31, 1.69], [after_rake, after_rake], color=COL_MUTED,
                linewidth=0.9, linestyle=":")
        ax.plot([2.31, 2.69], [banked, banked], color=COL_MUTED, linewidth=0.9, linestyle=":")

        ax.set_xticks(range(4), [s[0] for s in steps], fontsize=8.5)
        ax.axhline(0, color=COL_MUTED, linewidth=1.0)
        ax.set_title(name, fontsize=11, pad=8)
        ax.grid(axis="y", alpha=0.9)
        ax.set_axisbelow(True)

    axes[0][0].set_ylim(y_min, y_max)
    axes[0][0].set_ylabel(unit)
    fig.tight_layout()
    return fig


def _waterfall_slide(prs, layouts, config: Config, screens, in_euros: bool):
    # "in_euros" is the historical name for "in money rather than big blinds" -
    # the panel follows the display currency like everything else.
    money = config.currency.code
    title = (
        f"Where the money goes, in {money}" if in_euros
        else "Where the money goes, in big blinds"
    )
    slide = pc.add_image_slide(prs, layouts, title, _waterfall_figure(screens, config, in_euros))
    _footnote(
        slide,
        "Shared y axis across stakes - that comparison is the point. "
        + (
            f"In {money} every bar grows with the stake, which is the entire case for "
            "moving up."
            if in_euros
            else "In big blinds the rake bar shrinks as stakes rise, but the win rate shrinks "
            "faster: cheaper rake does not, on this data, pay for the tougher games."
        ),
    )
    return slide


# --------------------------------------------------------------------------- #
# Slide 5 - the chosen configuration and its neighbours
# --------------------------------------------------------------------------- #
def _configurations_slide(prs, layouts, config: Config, allocations, edge, best, current):
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "The chosen mix, and its nearest alternatives")

    # Two safer and two bolder neighbours ON THE FRONTIER, so every row shown is
    # a real option rather than an arbitrary nearby allocation.
    index = next((i for i, a in enumerate(edge) if best and a.counts == best.counts), None)
    if index is None:
        window = edge[:5]
        chosen = None
    else:
        low = max(0, index - 2)
        window = edge[low:index + 3]
        chosen = best.counts

    rows_data = [] if current is None else [_benchmark_row(current)]
    for allocation in window:
        is_best = chosen is not None and allocation.counts == chosen
        if is_best:
            note, bg = (
                f"CHOSEN - highest {config.currency.code}/hr inside tolerance", pc.COL_GREEN
            )
        elif allocation.risk_of_ruin < (best.risk_of_ruin if best else 0):
            note, bg = "safer, earns less", None
        else:
            # Only says "over your tolerance" when RUIN alone is what was tested,
            # since that is the column the reader can check against this row. Where
            # the downswing rule is in play a bolder row may be outside it for a
            # reason no column here shows, so the note stays descriptive.
            note, bg = (
                ("over your tolerance", pc.TABLE_LABEL_BG)
                if config.risk_mode == "ruin"
                else ("bolder than the chosen mix", pc.TABLE_LABEL_BG)
            )
        rows_data.append((allocation, note, bg, is_best))

    _allocation_table(slide, config, rows_data, pc.CONTENT_TOP + Inches(0.55))
    _footnote(
        slide,
        f"{risk_rule_note(config)} On a {_money(config, config.bankroll_eur)} bankroll "
        f"across {config.tables} tables. Risk of ruin assumes this mix is played at a fixed "
        "size forever, so it is an upper bound - in practice you would move down. "
        f"SD is the {config.currency.code} standard deviation per 100 hands: the risk being "
        "bought. " + drawdown_note(config),
    )
    return slide


def _allocation_table(slide, config: Config, rows_data, top):
    """The shared allocation table: configuration, money, risk, a note.

    Used by the frontier-neighbours slide and the single-stake slide so the two
    can be read against each other without the eye re-learning a layout.
    `rows_data` is a list of (allocation, note, background, bold).
    """
    code = config.currency.code
    columns = _fit_columns([
        ("Configuration", 2.55),
        (f"{code}/hr", 0.80),
        (f"On tables\n{code}", 0.95),
        (f"{timescale_label(config)} hands\nEV", 1.00),
        (f"Banked\n{code}/100", 0.95),
        (f"SD\n{code}/100", 0.90),
        ("Risk of\nruin", 0.85),
        # Three points on one distribution, so they share a stem and differ only
        # in the percentile - which is the comparison the reader wants to make.
        ("Downswing\nmedian", 1.00),
        ("Downswing\n10% worst", 1.00),
        ("Downswing\n1% worst", 1.00),
        ("", 1.70),
    ])
    table_width = Inches(sum(w for _, w in columns))
    left = int(pc.CONTENT_LEFT + (pc.CONTENT_WIDTH - table_width) / 2)
    shape = slide.shapes.add_table(
        len(rows_data) + 1, len(columns), left, top,
        table_width, Inches(0.45 + 0.38 * len(rows_data)),
    )
    table = shape.table
    for i, (_, width) in enumerate(columns):
        table.columns[i].width = Inches(width)
    for i, (label, _) in enumerate(columns):
        cell = table.cell(0, i)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=10, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE, wrap=True)

    for row, (allocation, note, bg, bold) in enumerate(rows_data, start=1):
        swing = sim.expected_drawdown(config, allocation, config.timescale_hands)
        values = [
            allocation.label,
            _plain(config, allocation.eur_per_hour),
            _plain(config, allocation.exposure_eur),
            _plain(config, horizon_ev(config, allocation)),
            _plain(config, allocation.mean_eur_per_100, 2),
            _plain(config, allocation.stdev_eur_per_100),
            f"{allocation.risk_of_ruin:.2%}",
            _plain(config, swing["median"]),
            _plain(config, swing["p90"]),
            _plain(config, swing["p99"]),
            note,
        ]
        for i, value in enumerate(values):
            cell = table.cell(row, i)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, value, font_size=10, bold=bold, bg_colour=bg)
    return table


def _single_stake_slide(prs, layouts, config: Config, screens, best, current):
    """Every stake played on its own, in the same format as the mix table.

    The baseline the whole exercise argues against: put all twelve tables on one
    stake and this is what you get. Shown immediately before the frontier so the
    reader has the single-stake numbers in mind when the mixes appear.
    """
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "If you played one stake and nothing else")

    rows_data = [] if current is None else [_benchmark_row(current)]
    for index, screen in enumerate(screens):
        counts = tuple(config.tables if i == index else 0 for i in range(len(config.stakes)))
        allocation = mix.evaluate(counts, config)
        if not screen.kept:
            note, bg = f"excluded - {screen.excluded_reason}", pc.TABLE_LABEL_BG
        elif allocation.within_ruin_tolerance:
            note, bg = "inside ruin tolerance", None
        else:
            note, bg = "over ruin tolerance", pc.TABLE_LABEL_BG
        rows_data.append((allocation, note, bg, False))

    _allocation_table(slide, config, rows_data, pc.CONTENT_TOP + Inches(0.45))
    best_line = (
        f" The best mix does {_money(config, best.eur_per_hour)}/hr at "
        f"{best.risk_of_ruin:.2%} - compare that against every row here."
        if best is not None else ""
    )
    _footnote(
        slide,
        f"All {config.tables} tables on one stake, the configuration conventional bankroll "
        "advice assumes. Note how quickly risk of ruin climbs relative to the earnings: the "
        "rungs are far apart in risk and close together in money, which is exactly the gap a "
        "mix exploits." + best_line + " " + drawdown_note(config),
    )
    return slide


# --------------------------------------------------------------------------- #
# The price of stepping up
# --------------------------------------------------------------------------- #
def _comparison_table(slide, config: Config, baseline, baseline_note, rows_data, top,
                      current=None):
    """Shared table: the played mix, the highlighted baseline, then variants.

    Deltas are always measured against the BASELINE (the optimum), including on
    the played-mix row - one reference point per table. The played mix is a
    benchmark to read against, not a second origin.
    """
    code = config.currency.code
    columns = _fit_columns([
        ("", 1.80),
        ("Configuration", 2.30),
        (f"{code}/hr", 0.75),
        ("vs optimal", 0.75),
        (f"On tables\n{code}", 0.90),
        (f"{timescale_label(config)} hands\nEV", 1.00),
        ("Risk of\nruin", 0.85),
        ("vs optimal", 0.75),
        ("Downswing\nmedian", 0.95),
        ("Downswing\n10% worst", 0.95),
        ("Downswing\n1% worst", 0.95),
        ("", 1.05),
    ])
    table_width = Inches(sum(w for _, w in columns))
    left = int(pc.CONTENT_LEFT + (pc.CONTENT_WIDTH - table_width) / 2)
    lead = 2 if current is not None else 1  # header, plus the benchmark row
    shape = slide.shapes.add_table(
        len(rows_data) + lead + 1, len(columns), left, top,
        table_width, Inches(0.45 + 0.40 * (len(rows_data) + lead)),
    )
    table = shape.table
    for i, (_, width) in enumerate(columns):
        table.columns[i].width = Inches(width)
    for i, (label, _) in enumerate(columns):
        cell = table.cell(0, i)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=10, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE, wrap=True)

    def swings(allocation):
        value = sim.expected_drawdown(config, allocation, config.timescale_hands)
        return [
            _plain(config, value["median"]),
            _plain(config, value["p90"]),
            _plain(config, value["p99"]),
        ]

    def delta(allocation):
        """A signed difference against the baseline, in display money."""
        gap = config.currency.from_eur(allocation.eur_per_hour - baseline.eur_per_hour)
        return f"{gap:+,.0f}"

    if current is not None:
        current_row = [
            CURRENT_LABEL, current.label, _plain(config, current.eur_per_hour),
            delta(current),
            _plain(config, current.exposure_eur),
            _plain(config, horizon_ev(config, current)),
            f"{current.risk_of_ruin:.2%}",
            f"{current.risk_of_ruin / max(baseline.risk_of_ruin, 1e-12):.1f}x",
            *swings(current), "what you play now",
        ]
        for i, value in enumerate(current_row):
            cell = table.cell(1, i)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, value, font_size=10, bold=True,
                        bg_colour=pc.COL_ORANGE, wrap=True)

    base_row = [
        baseline_note, baseline.label, _plain(config, baseline.eur_per_hour), "-",
        _plain(config, baseline.exposure_eur),
        _plain(config, horizon_ev(config, baseline)),
        f"{baseline.risk_of_ruin:.2%}", "-", *swings(baseline), "",
    ]
    for i, value in enumerate(base_row):
        cell = table.cell(lead, i)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, value, font_size=10, bold=True, bg_colour=pc.COL_GREEN, wrap=True)

    for row, (label, allocation, note, inside) in enumerate(rows_data, start=lead + 1):
        values = [
            label, allocation.label, _plain(config, allocation.eur_per_hour),
            delta(allocation),
            _plain(config, allocation.exposure_eur),
            _plain(config, horizon_ev(config, allocation)),
            f"{allocation.risk_of_ruin:.2%}",
            f"{allocation.risk_of_ruin / max(baseline.risk_of_ruin, 1e-12):.1f}x",
            *swings(allocation), note,
        ]
        for i, value in enumerate(values):
            cell = table.cell(row, i)
            pc._zero_cell_margins(cell)
            pc.set_cell(
                cell, value, font_size=10, wrap=True,
                bg_colour=None if inside else pc.TABLE_LABEL_BG,
                # The note column, wherever it has ended up - an index literal
                # here is what breaks silently when a column is inserted.
                font_colour=(
                    pc.COL_TEXT_RED if (not inside and i == len(values) - 1) else None
                ),
            )
    return table


def _step_up_slide(prs, layouts, config: Config, best, options, current):
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "The two ways up, and what each costs")

    if not options:
        box = slide.shapes.add_textbox(
            pc.CONTENT_LEFT, pc.CONTENT_TOP + Inches(0.6), pc.CONTENT_WIDTH, Inches(1)
        )
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "The optimal mix already sits at the top of the ladder."
        run.font.size = Pt(12)
        return slide

    rows = [
        (
            option.label,
            option.allocation,
            "inside tolerance" if option.within_tolerance else "OVER tolerance",
            option.within_tolerance,
        )
        for option in options
    ]
    _comparison_table(slide, config, best, "OPTIMAL (baseline)", rows,
                      pc.CONTENT_TOP + Inches(0.55), current=current)
    used = [i for i, count in enumerate(best.counts) if count]
    shot = config.stakes[max(used) + 1].name
    highest, lowest = config.stakes[max(used)].name, config.stakes[min(used)].name
    _footnote(
        slide,
        f"Both moves take the SAME shot - one table on {shot}, the rung above the top of the "
        f"optimal mix - and both keep the table count fixed. They differ only in which table "
        f"pays for it: TOP UP gives up a {highest} table, stretching the top of the mix up a "
        f"rung; BOTTOM UP gives up a {lowest} table, leaving the top where it is and thinning "
        f"the bottom instead. " + drawdown_note(config),
    )
    return slide


BANKROLL_STEPS = (0.0, 0.10, 0.25, 0.50, 1.00)
"""Growth multiples to re-solve the optimum at. Deliberately proportional rather
than absolute: what matters is how far the roll has moved, not where it started."""


def bankroll_ladder(config: Config):
    """The optimum re-solved at each growth step.

    Yields (label, bankroll, scenario config, allocation). Each row is a fresh,
    independent problem - the same closed-form calculation the deck opens with,
    run at a larger roll. Shared with the simulation-priming pass so the
    downswing figures are computed once, at the RIGHT bankroll for each row.
    """
    rows = []
    for step in BANKROLL_STEPS:
        bankroll = config.bankroll_eur * (1 + step)
        scenario = config.replace(bankroll_eur=bankroll)
        # Re-solved under the LIVE risk rule. In downswing mode each step is its
        # own walk, and the drawdown cache is keyed on bankroll (correctly - the
        # absorbing barrier moves with it), so the steps do not share simulations.
        allocation = mix.best_allocation(mix.all_allocations(scenario), scenario)
        rows.append((f"+{step:.0%}" if step else "now", bankroll, scenario, allocation))
    return rows


def _bankroll_ladder_slide(prs, layouts, config: Config, best, current):
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "What to play as the bankroll grows")

    code = config.currency.code
    columns = _fit_columns([
        ("Bankroll", 0.80),
        ("", 1.05),
        ("Optimal mix", 2.45),
        (f"{code}/hr", 0.75),
        ("vs now", 0.75),
        (f"On tables\n{code}", 0.90),
        (f"{timescale_label(config)} hands\nEV", 1.00),
        ("Risk of\nruin", 0.85),
        ("Downswing\nmedian", 0.95),
        ("Downswing\n10% worst", 0.95),
        ("Downswing\n1% worst", 0.95),
        ("", 1.05),
    ])
    rows = bankroll_ladder(config)
    lead = 2 if current is not None else 1  # header, plus the benchmark row
    table_width = Inches(sum(w for _, w in columns))
    left = int(pc.CONTENT_LEFT + (pc.CONTENT_WIDTH - table_width) / 2)
    shape = slide.shapes.add_table(
        len(rows) + lead, len(columns), left, pc.CONTENT_TOP + Inches(0.55),
        table_width, Inches(0.45 + 0.40 * (len(rows) + lead - 1)),
    )
    table = shape.table
    for i, (_, width) in enumerate(columns):
        table.columns[i].width = Inches(width)
    for i, (label, _) in enumerate(columns):
        cell = table.cell(0, i)
        pc._zero_cell_margins(cell)
        pc.set_cell(cell, label, font_size=10, bold=True,
                    bg_colour=pc.TABLE_HEADER_BG, font_colour=pc.WHITE, wrap=True)

    if current is not None:
        # Priced at TODAY's bankroll, because that is the roll it is being
        # played on - the rows below re-solve at bigger rolls, this one does not.
        swing = sim.expected_drawdown(config, current, config.timescale_hands)
        values = [
            "playing", _money(config, config.bankroll_eur), current.label,
            _plain(config, current.eur_per_hour),
            f"{config.currency.from_eur(current.eur_per_hour - best.eur_per_hour):+,.0f}"
            if best else "-",
            _plain(config, current.exposure_eur),
            _plain(config, horizon_ev(config, current)),
            f"{current.risk_of_ruin:.2%}",
            _plain(config, swing["median"]), _plain(config, swing["p90"]),
            _plain(config, swing["p99"]),
            "what you play now",
        ]
        for i, value in enumerate(values):
            cell = table.cell(1, i)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, value, font_size=10, bold=True,
                        bg_colour=pc.COL_ORANGE, wrap=True)

    previous_top = None
    for row, (label, bankroll, scenario, allocation) in enumerate(rows, start=lead):
        is_now = label == "now"
        if allocation is None:
            values = [label, _money(config, bankroll), "nothing clears tolerance",
                      "-", "-", "-", "-", "-", "-", "-", "-", ""]
        else:
            top = max(i for i, c in enumerate(allocation.counts) if c)
            note = "where you should be now" if is_now else (
                f"unlocks {config.stakes[top].name}"
                if previous_top is not None and top > previous_top else ""
            )
            previous_top = top
            swing = sim.expected_drawdown(scenario, allocation, config.timescale_hands)
            values = [
                label,
                _money(config, bankroll),
                allocation.label,
                _plain(config, allocation.eur_per_hour),
                f"{config.currency.from_eur(allocation.eur_per_hour - best.eur_per_hour):+,.0f}"
                if best else "-",
                _plain(config, allocation.exposure_eur),
                # Priced on the ROW's own scenario, so the EV of a bigger roll
                # is the EV of the mix it can actually afford.
                _plain(config, horizon_ev(scenario, allocation)),
                f"{allocation.risk_of_ruin:.2%}",
                _plain(config, swing["median"]),
                _plain(config, swing["p90"]),
                _plain(config, swing["p99"]),
                note,
            ]
        for i, value in enumerate(values):
            cell = table.cell(row, i)
            pc._zero_cell_margins(cell)
            pc.set_cell(cell, value, font_size=10, bold=is_now,
                        bg_colour=pc.COL_GREEN if is_now else None, wrap=True)

    _footnote(
        slide,
        "Each row is the optimum re-solved from scratch at that bankroll, under the same "
        f"rule - a fresh problem, not a path. {risk_rule_note(config)} It does not "
        "model getting there: no move-down rule and no order of play, just the right answer "
        "if you were standing at that roll today. Risk of ruin and both downswing figures are "
        "measured from that row's bankroll as the starting point. " + drawdown_note(config),
    )
    return slide


# --------------------------------------------------------------------------- #
# Slide 7 - the simulation
# --------------------------------------------------------------------------- #
def _simulation_slide(prs, layouts, config: Config, result, title, scales=(None, None)):
    import numpy as np

    ylim, drawdown_xmax = scales
    slide = pc.add_image_slide(
        prs, layouts, title,
        charts.simulation_figure(result, config, ylim=ylim, drawdown_xmax=drawdown_xmax),
    )
    median_dd = float(np.percentile(result.max_drawdown, 50))
    p90_dd = float(np.percentile(result.max_drawdown, 90))
    median_end = float(np.percentile(result.final_bankroll, 50))
    over_roll = (
        " Note the typical worst drawdown EXCEEDS the starting bankroll while only "
        f"{result.ruin_probability:.2%} go broke: by the time the deep falls arrive, most "
        "lifetimes are playing with winnings, not with the money they started on."
        if median_dd > config.bankroll_eur
        else ""
    )
    _footnote(
        slide,
        f"{result.allocation.label}. {result.hands:,} hands ({result.hours:,.0f} hours at "
        f"{config.tables} tables) per lifetime, played at a fixed mix with no move-down rule - "
        f"an upper bound on risk, not a forecast. Typical worst drawdown "
        f"{_money(config, median_dd)}, one lifetime in ten worse than "
        f"{_money(config, p90_dd)}, median finish {_money(config, median_end)}." + over_roll,
    )
    return slide


def _random_paths_slide(prs, layouts, config: Config, result, title, best, current,
                        scales=(None, None)):
    """The fan chart's companion: twenty individual lifetimes, not bands.

    Both EV lines appear on BOTH copies of this slide, so the comparison is the
    same one either way round - the question is always 'where does this mix end
    up against the other one', and having the reference line move between slides
    would make the two impossible to read together."""
    import numpy as np

    ev_lines = []
    if best is not None:
        ev_lines.append(("Optimal", best, charts.STATUS_GOOD))
    if current is not None and (best is None or current.counts != best.counts):
        # Pink, not red: red is the ruin barrier on this very chart, and the mix
        # you are playing is a reference line, not a hazard.
        ev_lines.append(("Current", current, charts.COL_CURRENT))

    slide = pc.add_image_slide(
        prs, layouts, title,
        charts.random_paths_figure(result, config, ev_lines, ylim=scales[0]),
    )

    hands = result.hands
    spread = (
        f"{_money(config, float(np.percentile(result.final_bankroll, 5)))} to "
        f"{_money(config, float(np.percentile(result.final_bankroll, 95)))}"
    )
    ev_note = "  ".join(
        f"{label} EV finishes at "
        f"{_money(config, config.bankroll_eur + allocation.mean_eur_per_100 * hands / 100)}."
        for label, allocation, _ in ev_lines
    )
    _footnote(
        slide,
        f"{result.allocation.label}. Twenty of the {result.paths:,} simulated lifetimes, "
        f"drawn at random rather than picked by where they finish, so the spread is the one "
        f"the simulation produced. Ninety per cent of lifetimes finish between {spread}, which "
        f"is the width the EV lines do not show. {ev_note} The EV lines are straight because "
        f"expectation is linear in hands - nothing here compounds, since the mix is static and "
        f"the stakes never move. " + drawdown_note(config),
    )
    return slide


# --------------------------------------------------------------------------- #
# Appendix - why "more money, less risk" is not a contradiction
# --------------------------------------------------------------------------- #
def _dominance_slide(prs, layouts, config: Config, allocations, best, current):
    """Why the chart has mixes that beat others on BOTH axes.

    The single most-asked question about this deck, and the intuition it breaks
    is a reasonable one: surely more money always costs more risk. Every figure
    on the slide is recomputed from the same functions the tables use.
    """
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "Appendix: how a mix can earn MORE and risk LESS")

    stakes = config.stakes
    low, high = stakes[0], stakes[-1]

    def solo(stake):
        index = stakes.index(stake)
        counts = tuple(config.tables if i == index else 0 for i in range(len(stakes)))
        return mix.evaluate(counts, config)

    solo_low, solo_high = solo(low), solo(high)
    ev_multiple = solo_high.mean_eur_per_100 / solo_low.mean_eur_per_100
    var_multiple = solo_high.variance_eur_per_100 / solo_low.variance_eur_per_100

    dominated = [
        a for a in allocations
        if any(b.eur_per_hour > a.eur_per_hour and b.risk_of_ruin < a.risk_of_ruin
               for b in allocations)
    ]

    sections = [
        ("Risk of ruin is not the size of the swings", [
            "It is the size of the swings measured against how fast you win. Money coming "
            "in is what pulls the bankroll away from zero, so your edge is your protection.",
            "Two players with identical swings, one winning twice as fast: the faster one is "
            "far safer. They climb out of holes before the holes get deep.",
        ]),
        ("Which is why small stakes are not safe - they are slow", [
            f"You have {config.tables} seats and that is the scarce thing. A seat at "
            f"{low.name} adds almost no swing, but it adds almost no income either. It does "
            "not make you safer, it dilutes your win rate across the whole operation.",
            "A diluted win rate with unchanged swings is worse on the only ratio that "
            "matters. Slow is what actually breaks you.",
        ]),
        ("Moving UP a rung, though, is never free", [
            f"Earnings scale with the big blind; variance scales with its SQUARE. All "
            f"{config.tables} tables on {high.name} rather than {low.name} is "
            f"{ev_multiple:.1f}x the earnings and {var_multiple:.0f}x the variance.",
            f"So ruin climbs the whole way up the ladder - {solo_low.risk_of_ruin:.2%} on "
            f"{low.name} against {solo_high.risk_of_ruin:.1%} on {high.name}. Moving up always "
            "buys risk with money. That is a real trade-off and it is what the frontier shows.",
        ]),
        ("So where does 'more money, less ruin' come from?", [
            "Not from moving up. From moving seats that earn nothing onto stakes that do: the "
            "swings grow, but the earnings grow FASTER, so the ratio improves and ruin falls.",
            f"That is why {len(dominated):,} of the {len(allocations):,} possible mixes "
            f"({len(dominated) / len(allocations):.0%}) are beaten by another mix on both "
            "axes at once. They are not taking a different trade-off - they are simply wasting "
            "risk they have already accepted.",
        ]),
    ]

    if current is not None and best is not None:
        beaten = [
            a for a in allocations
            if a.eur_per_hour > current.eur_per_hour
            and a.risk_of_ruin < current.risk_of_ruin
        ]
        worst_seat = min(
            (i for i, count in enumerate(current.counts) if count), default=None
        )
        detail = [
            f"The mix in play earns {_money(config, current.eur_per_hour)}/hr at "
            f"{current.risk_of_ruin:.2%} ruin. {len(beaten):,} other mixes earn MORE and risk "
            "LESS - it is one of the wasteful ones, not a cautious one."
            if beaten else
            f"The mix in play earns {_money(config, current.eur_per_hour)}/hr at "
            f"{current.risk_of_ruin:.2%} ruin, and nothing beats it on both axes at once.",
        ]
        if worst_seat is not None:
            detail.append(
                f"The {current.counts[worst_seat]} seats at {stakes[worst_seat].name} are the "
                f"drag, not the top of the mix: the optimum plays "
                f"{best.counts[-1] if best.counts[-1] else max(best.counts)} tables at its own "
                f"top rung. Cutting the highest stake is the wrong instinct."
            )
        detail.append(
            f"And the tolerance is {config.ruin_tolerance:.2%} while the mix in play runs at "
            f"{current.risk_of_ruin:.2%} - roughly "
            f"{config.ruin_tolerance / max(current.risk_of_ruin, 1e-12):.0f}x less risk than "
            f"has been authorised. That unused allowance is worth "
            f"{config.currency.from_eur(best.eur_per_hour - current.eur_per_hour):+,.0f} "
            f"{config.currency.code}/hr."
        )
        sections.append(("Your case, in three numbers", detail))

    body = slide.shapes.add_textbox(
        pc.CONTENT_LEFT, pc.CONTENT_TOP, pc.CONTENT_WIDTH, pc.CONTENT_HEIGHT
    )
    frame = body.text_frame
    frame.word_wrap = True
    first = True
    for heading, points in sections:
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        para.space_before = Pt(0 if first else 8)
        first = False
        run = para.add_run()
        run.text = heading
        run.font.size = Pt(12)
        run.font.bold = True
        for point in points:
            bullet = frame.add_paragraph()
            bullet.level = 1
            run = bullet.add_run()
            run.text = point
            run.font.size = Pt(10)
            run.font.color.rgb = pc.BLACK
    return slide


# --------------------------------------------------------------------------- #
# Slide 8 - methodology
# --------------------------------------------------------------------------- #
def _methodology_slide(prs, layouts, config: Config):
    slide = prs.slides.add_slide(layouts["Title and Content"])
    pc.add_title(slide, "Appendix: method, and what it assumes")

    body = slide.shapes.add_textbox(
        pc.CONTENT_LEFT, pc.CONTENT_TOP, pc.CONTENT_WIDTH, pc.CONTENT_HEIGHT
    )
    frame = body.text_frame
    frame.word_wrap = True

    sections = [
        ("The question", [
            "Which distribution of tables across stakes earns the most, without exceeding a "
            "chosen risk tolerance. Volume is divisible across tables, so the decision variable "
            "is the share of tables at each stake, not a single stake.",
            risk_rule_note(config),
        ]),
        ("The maths", [
            "Per 100 hands dealt across all tables, a mix has mean  SUM (n/T) x winrate x bb, "
            "and variance  SUM (n/T) x sd^2 x bb^2, both in euros. Tables deal independent "
            "hands, so these are exact, not simulated.",
            "Risk of ruin follows the standard diffusion result  R = exp(-2 x mean x bankroll "
            "/ variance).  P(lose half) is the same formula with the barrier moved up.",
            "Every allocation is enumerated and scored - a few hundred to a few thousand - so "
            "the optimum is found by brute force. No sampling and no search heuristic.",
        ]),
        ("Rakeback", [
            f"Rebated at {config.rakeback_pct:.0%} of rake paid. It is a rebate on volume, not "
            "a gamble: it raises the mean and adds NO variance, which is why it both increases "
            "earnings and reduces risk of ruin.",
            "Win rates are already net of rake, so adding rakeback does not double-count.",
        ]),
        ("What is assumed, not measured", [
            # Was "a flat 92 bb/100 at every stake", which the config has not
            # been for some time - it is measured wherever the sample allows.
            "Standard deviation is MEASURED at the stakes with a real sample and ASSUMED at "
            f"92 bb/100 at the thin ones - see the run-parameters slide for which is which. "
            "A measured SD off a few hundred hands would understate the variance badly, which "
            "is why those are overridden rather than believed.",
            "Tables are treated as independent. Multi-tabling correlates outcomes within a "
            "session, which would raise effective variance.",
            "No win-rate penalty is charged for table count.",
            "Ruin assumes the mix is played at a fixed size forever, with no move-down rule. "
            "That makes it a conservative constraint and a poor forecast.",
        ]),
        ("Known weaknesses in the inputs", [
            "Only the lowest stakes are well measured. A win rate's 95% interval is roughly "
            "+/-3.5 bb/100 at 271k hands and +/-16.5 at 12k, so the higher rungs of the ladder "
            "rest on thin samples.",
            "Stakes with a few hundred hands are held out entirely rather than allowed to "
            "dominate the answer on noise.",
        ]),
        ("Two kinds of downswing, and why only one has a closed form", [
            "LOSS BELOW STARTING BANKROLL is what ruin measures - start on 10k, run to 15k, "
            "fall to 5k, and that is a 5k loss below start. Over unlimited time it is bounded "
            "and exactly exponential, so its quantiles are exact. The 99th percentile of it "
            "IS your bankroll, by construction, at a 1% tolerance.",
            "PEAK-TO-TROUGH DRAWDOWN is what a downswing feels like - the same episode above "
            "is a 10k fall. It has NO all-time value: given unlimited time it grows without "
            "bound, because a winning bankroll keeps making new highs to fall from. Every "
            "peak-to-trough figure here is therefore 'within this many hands' and nothing "
            "more. That is the reason the simulation exists.",
        ]),
        ("The simulation", [
            f"{config.timescale_hands:,} hands per lifetime, {config.sim_paths:,} independent "
            "lifetimes, ruin absorbing (a busted path stops rather than trading back).",
            "Hands are drawn in 100-hand blocks, the unit the win rate and variance arrive "
            "in - exact for a diffusion, though a dip and recovery inside one block goes "
            "unseen.",
            "Checked against the closed form: over a long horizon its loss-below-start "
            "distribution must reproduce the analytic exponential. A simulation that "
            "disagrees where the maths is known has a bug.",
        ]),
        ("Where the CURRENT row on every table comes from", [
            "Hands actually played per stake in the period reviewed, apportioned to whole "
            "tables by largest remainder. It is a reconstruction of how the table time was "
            "really split, not a mix that was consciously chosen.",
            "Only the RATIO between stakes matters, so the absolute hand counts and the exact "
            "period they cover do not affect it.",
        ]),
        ("Not modelled", [
            "The allocation is a static snapshot at one bankroll, in the analytics AND in the "
            "simulation. The dynamic version - move up through one threshold, down through "
            "another, with hysteresis - is not built. Since a real player moves down long "
            "before busting, every risk number here is an upper bound.",
        ]),
    ]

    note = config.currency.note()
    if note:
        sections.append((
            "Currency",
            [
                note,
                "The rate is a fixed constant, not a live quote - a bankroll plan that moved "
                "with spot would not be a plan, and rebuilding the deck would silently change "
                "every figure on it. The tables themselves are dealt in EUR throughout, and "
                "the big blinds on the run-parameters slide stay in EUR for that reason.",
            ],
        ))

    first = True
    for heading, points in sections:
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        run = para.add_run()
        run.text = heading
        run.font.size = Pt(12)
        run.font.bold = True
        para.space_before = Pt(0 if heading == "The question" else 7)
        for point in points:
            bullet = frame.add_paragraph()
            bullet.level = 1
            run = bullet.add_run()
            run.text = point
            run.font.size = Pt(10)
            run.font.color.rgb = pc.BLACK
    return slide


# --------------------------------------------------------------------------- #
def _footnote(slide, text):
    box = slide.shapes.add_textbox(
        pc.CONTENT_LEFT,
        pc.CONTENT_TOP + pc.CONTENT_HEIGHT - Inches(0.85),
        pc.CONTENT_WIDTH,
        Inches(0.8),
    )
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = pc.GRID_GREY
    return box


def _mixes_needing_simulation(config: Config, screens, edge, best, current):
    """Every allocation that will appear in a table, de-duplicated, in order."""
    wanted = []

    def want(allocation):
        if allocation is not None and allocation.counts not in {a.counts for a in wanted}:
            wanted.append(allocation)

    for index in range(len(config.stakes)):  # the single-stake slide
        want(mix.evaluate(
            tuple(config.tables if i == index else 0 for i in range(len(config.stakes))),
            config,
        ))
    if best is not None:
        position = next((i for i, a in enumerate(edge) if a.counts == best.counts), None)
        window = edge[max(0, position - 2):position + 3] if position is not None else edge[:5]
        for allocation in window:
            want(allocation)
        for option in mix.step_up_options(config, best):
            want(option.allocation)
    want(current)
    return wanted


def _prime_simulations(config: Config, screens, edge, best, current) -> dict:
    """Run every simulation the deck needs, with a progress bar.

    Returns the full-detail results for the headline charts, keyed by
    allocation, so they are not simulated a second time when those slides are
    built. The per-table figures go into sim's own cache.
    """
    wanted = _mixes_needing_simulation(config, screens, edge, best, current)
    headline = [a for a in (best, current) if a is not None]
    # The ladder rows each carry their OWN bankroll, so they cannot share the
    # cache entries above even where the allocation happens to match.
    ladder = [
        (scenario, allocation, label)
        for label, _, scenario, allocation in bankroll_ladder(config)
        if allocation is not None and label != "now"
    ]

    total = len(wanted) + len(ladder) + len(headline)
    print(f"\nSTEP 3 - SIMULATIONS   ({total} mixes x {config.timescale_hands:,} hands)")
    bar = progress.Progress(total, "simulations")
    bar.draw("starting")

    for allocation in wanted:
        sim.expected_drawdown(config, allocation, config.timescale_hands)
        bar.advance(allocation.label)

    for scenario, allocation, label in ladder:
        sim.expected_drawdown(scenario, allocation, config.timescale_hands)
        bar.advance(f"{allocation.label} (bankroll {label})")

    results = {}
    for allocation in headline:
        # The headline charts run at the full path count for a smooth histogram;
        # the table figures use fewer, which is plenty for a percentile.
        results[allocation.counts] = sim.simulate(
            config, allocation, hands=config.timescale_hands, paths=config.sim_paths
        )
        bar.advance(f"{allocation.label} (full detail)")

    bar.close(f"{total} mixes simulated over {config.timescale_hands:,} hands each.")
    return results


def build(config: Config, directory: Path, workbook_path: Path | None = None) -> Path:
    """Build the deck and return the path written.

    With `workbook_path`, the same numbers are also written there as a
    spreadsheet. It happens INSIDE this function rather than beside it because
    every value the workbook wants is one this function already has: doing it
    from outside would mean a second tolerance walk and a second round of
    simulations, and any drift between the two would ship as a spreadsheet that
    contradicts its own deck.

    The caller supplies the path, so it knows what was written without this
    returning a second value and disturbing every existing call site.
    """
    directory.mkdir(parents=True, exist_ok=True)
    path = directory / "stake_optimisation.pptx"

    screens = mix.screen_stakes(config)
    allocations = mix.all_allocations(config)
    edge = mix.frontier(allocations)
    best = mix.best_allocation(
        allocations,
        config,
        progress_label="testing mixes" if config.risk_mode != "ruin" else None,
    )

    current = mix.current_allocation(config)

    # Simulate everything the deck will need BEFORE building any slides, so the
    # work is visible as one progress bar rather than as a series of unexplained
    # pauses while tables render. Results land in sim's cache; the table code
    # then finds them already there.
    simulations = _prime_simulations(config, screens, edge, best, current)

    prs, layouts = pc.load_template_presentation()

    def run_sim(allocation):
        return simulations[allocation.counts]

    # ONE set of axis limits for every simulation chart, computed across all of
    # them before any is drawn. Two mixes autoscaled separately look far more
    # alike than they are, and these slides exist to be compared by eye.
    plotted = [a for a in (best, current) if a is not None]
    scales = (
        charts.simulation_scales([run_sim(a) for a in plotted], config, plotted)
        if plotted else (None, None)
    )

    # ---- 0. What this was run on ------------------------------------------- #
    _run_parameters_slide(prs, layouts, config)

    # ---- 1. Where the tables should be ------------------------------------- #
    pc.add_chapter_slide(
        prs, layouts, "1. Optimal stake distribution",
        "Which stakes are worth playing, and in what proportion",
    )
    _stake_table_slide(prs, layouts, config, screens)
    _winrate_ci_slide(prs, layouts, config, screens)
    _waterfall_slide(prs, layouts, config, screens, in_euros=False)
    _waterfall_slide(prs, layouts, config, screens, in_euros=True)
    _single_stake_slide(prs, layouts, config, screens, best, current)
    pc.add_image_slide(
        prs, layouts, "Every way to split the tables",
        charts.allocation_frontier_figure(config),
    )
    # The same trade-off on the other risk axis. Every point it needs has already
    # been simulated by this stage, so it is cheap here even though the chart is
    # the expensive one to build cold.
    pc.add_image_slide(
        prs, layouts, "The same split, priced in downswings",
        charts.allocation_frontier_downswing_figure(config),
    )
    _configurations_slide(prs, layouts, config, allocations, edge, best, current)
    if best is not None:
        _simulation_slide(
            prs, layouts, config, run_sim(best),
            "The optimal mix, simulated", scales,
        )
        _random_paths_slide(
            prs, layouts, config, run_sim(best),
            "The optimal mix, twenty single lifetimes", best, current, scales,
        )

    # ---- 2. Moving up ------------------------------------------------------ #
    if best is not None:
        pc.add_chapter_slide(
            prs, layouts, "2. Shot-taking",
            "What it costs to move a table up a rung",
        )
        _step_up_slide(prs, layouts, config, best, mix.step_up_options(config, best), current)

    # ---- 3. The same question at a bigger roll ----------------------------- #
    if best is not None:
        pc.add_chapter_slide(
            prs, layouts, "3. As the bankroll grows",
            "Where the mix should go as the roll gets bigger",
        )
        _bankroll_ladder_slide(prs, layouts, config, best, current)

    # ---- 4. What is actually happening ------------------------------------- #
    # No table slide here: the played mix is the first row of every table in the
    # deck, so a slide showing it once more would be the fourth copy. What only
    # this section can give is the two simulations of it.
    if current is not None and best is not None:
        pc.add_chapter_slide(
            prs, layouts, "4. Current configuration",
            "The mix actually played, simulated the same way as the optimum",
        )
        _simulation_slide(
            prs, layouts, config, run_sim(current),
            "The mix you actually played, simulated", scales,
        )
        _random_paths_slide(
            prs, layouts, config, run_sim(current),
            "The mix you actually played, twenty single lifetimes", best, current, scales,
        )

    _dominance_slide(prs, layouts, config, allocations, best, current)
    _methodology_slide(prs, layouts, config)

    prs.save(path)

    if workbook_path is not None:
        from . import workbook

        # Every simulated figure it asks for is already in sim's cache, so this
        # costs a few seconds of writing rather than another round of Monte Carlo.
        workbook.write(
            workbook_path, config, screens, allocations, edge, best, current,
            simulations,
            bankroll_ladder(config) if best is not None else [],
            mix.step_up_options(config, best) if best is not None else [],
        )

    return path

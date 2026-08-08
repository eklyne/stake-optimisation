"""
shotopt/pptx_common.py
----------------------
PPTX-building infrastructure: template loading, RGBColor constants, low-level
XML helpers, and placeholder-aware title/chapter/image helpers.

COPIED from the nemesis-mvp analytics repo (lib/pptx_common.py) rather than
imported, deliberately: this repo depends on nothing over there, and a shared
import would create exactly that dependency. The usual cost of a copy applies -
fixes do not flow between the two, so treat this as a fork taken on 2026-08-08,
not as a live mirror.

Two local changes from the original:

* TEMPLATE_PATH resolves to this repo's own assets/ directory.
* load_template_presentation() falls back to a blank 16:9 Presentation when the
  branded template is missing. The template is a personal asset and gitignored,
  so a fresh clone must still be able to build a plainer deck rather than
  crashing on an absent file.

Content-area geometry is read off the template's own layout placeholders, not
scaled from a 10x7.5in canvas.
"""

import os

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from xml.sax.saxutils import escape

# Resolved off this file's location (shotopt/ -> repo root) rather than the cwd,
# so a deck build launched from anywhere still finds the template.
TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "assets", "deck_template.pptx",
)

# ── Canvas / content-area geometry (matches the template's own placeholders) ──

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TITLE_LEFT   = Inches(0.425)
TITLE_TOP    = Inches(0.399)
TITLE_WIDTH  = Inches(12.483)
TITLE_HEIGHT = Inches(0.999)

CONTENT_LEFT   = Inches(0.425)
CONTENT_TOP    = Inches(1.495)
CONTENT_WIDTH  = Inches(12.483)
# Bottom margin stops at 7.3in (0.2in above the true 7.5in slide edge), not at
# the footer placeholder's 6.464in top -- an untouched placeholder (we never
# set text into the footer/date/slide-number placeholders) clones onto the
# slide as an empty text frame and renders nothing, so our own shapes can
# safely extend underneath it.
CONTENT_HEIGHT = Inches(5.805)

# ── Colour constants (kept hardcoded per user preference -- not derived from
# the template's theme) -- one definition shared by both deck-building scripts
# instead of two independently-maintained copies. ──────────────────────────────

WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
BLACK           = RGBColor(0x00, 0x00, 0x00)
COL_GREEN       = RGBColor(0x74, 0xC4, 0x76)
COL_RED         = RGBColor(0xFC, 0x92, 0x72)
# Bar fills above are deliberately soft; TEXT needs a strong red to read as a
# callout at slide distance (the soft one renders as orange on a white slide).
COL_TEXT_RED    = RGBColor(0xC0, 0x00, 0x00)
TABLE_HEADER_BG = RGBColor(0x2F, 0x2F, 0x2F)
TABLE_LABEL_BG  = RGBColor(0xF2, 0xF2, 0xF2)
CHAPTER_BG      = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE     = RGBColor(0x6B, 0xAE, 0xD6)
GRID_HDR_BG     = RGBColor(0x3C, 0x3C, 0x4C)
GRID_GREY       = RGBColor(0xC0, 0xC0, 0xC0)  # insufficient sample
GRID_NA         = RGBColor(0xE8, 0xE8, 0xE8)  # no data at all

BORDER_W_EMU = int(0.75 * 12700)


def template_available():
    """True when the branded template is on disk. It is gitignored, so a fresh
    clone has to cope without it."""
    return os.path.isfile(TEMPLATE_PATH)


def load_template_presentation():
    """Open the template and strip every existing (filled-in) slide, keeping
    the master/layouts/theme. Returns (prs, layouts), where layouts maps
    layout name -> slide layout object.

    Without the template, falls back to python-pptx's default 4:3 deck widened
    to 16:9. The layout NAMES differ there, so `layouts` is padded with aliases
    for the two names the deck builders ask for by name ('Title and Content',
    '2_Section Divider') - a deck built this way is plain, but it builds."""
    if not template_available():
        prs = Presentation()
        prs.slide_width, prs.slide_height = SLIDE_WIDTH, SLIDE_HEIGHT
        layouts = {}
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                layouts.setdefault(layout.name, layout)
        # 'Title and Content' happens to exist in the default master; the
        # section divider does not, so it borrows the title-only layout.
        layouts.setdefault('Title and Content', prs.slide_layouts[1])
        layouts.setdefault('2_Section Divider', prs.slide_layouts[5])
        return prs, layouts

    prs = Presentation(TEMPLATE_PATH)

    # python-pptx has no high-level "delete slide" API -- drop each slide's
    # relationship and remove its <p:sldId> entry from the slide id list.
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)

    layouts = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layouts.setdefault(layout.name, layout)
    return prs, layouts


# ── XML helpers (ported unchanged from 11_plot_charts.py/12_seat_selection.py) ─

def _no_line(shape):
    """Remove visible border from an autoshape."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn('a:ln')):
        spPr.remove(old)
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))


def add_border_to_spPr(spPr):
    for old in spPr.findall(qn('a:ln')):
        spPr.remove(old)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(BORDER_W_EMU))
    etree.SubElement(etree.SubElement(ln, qn('a:solidFill')), qn('a:srgbClr')).set('val', str(BLACK))


def add_series_borders(chart):
    for series in chart.series:
        ser  = series._element
        spPr = ser.find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(ser, qn('c:spPr'))
        add_border_to_spPr(spPr)


def _zero_cell_margins(cell):
    """Zero a table cell's internal margins and vertically center its text --
    removes python-pptx's default cell padding so dense tables aren't clipped."""
    tcPr = cell._tc.get_or_add_tcPr()
    for attr in ('marL', 'marR', 'marT', 'marB'):
        tcPr.set(attr, '0')
    tcPr.set('anchor', 'ctr')


def set_cell(cell, text, font_size=9, bold=False, bg_colour=None, font_colour=None, wrap=False):
    """Low-level table-cell text setter: splits text on '\\n' into centered
    paragraphs, sets font size/bold/colour per run, and optionally fills the
    cell background via raw XML (python-pptx has no fill-by-hex shortcut used
    here)."""
    tf = cell.text_frame
    tf.word_wrap = wrap
    lines = str(text).split('\n') if text else ['']
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if font_colour:
            run.font.color.rgb = font_colour
    if bg_colour:
        tcPr = cell._tc.get_or_add_tcPr()
        etree.SubElement(etree.SubElement(tcPr, qn('a:solidFill')),
                         qn('a:srgbClr')).set('val', str(bg_colour))


# ── Placeholder-aware title/chapter helpers ────────────────────────────────────

def add_title(slide, text):
    """Set the slide's real title placeholder text -- font/position/colour
    all come from the template's layout instead of a manual textbox."""
    slide.shapes.title.text = text


def _to_rgb(colour):
    """Accept either an RGBColor or a hex string ('4C72B0' or '#4C72B0') and
    return an RGBColor -- lets callers pass either the pipeline's existing
    RGBColor constants (COL_GREEN/COL_RED/...) or a plain hex string (e.g.
    SEGMENT_COLOR_MAP's values) interchangeably."""
    if isinstance(colour, RGBColor):
        return colour
    return RGBColor.from_string(str(colour).lstrip('#'))


# ── Native (editable) chart helpers -- additive, for 14_line_optimisation.py's
# matplotlib-PNG-to-native-chart conversion. Deliberately NOT used to refactor
# 11_plot_charts.py's own add_segment_slide/add_winrate_slide/colour_bars_by_sign
# /add_pct_labels, which keep their own independent copies per this module's
# own "duplicate rather than couple two top-level scripts" precedent -- these
# are new, generic building blocks for a THIRD caller (14) rather than a
# refactor of the first two. ──────────────────────────────────────────────────

def add_bar_chart(slide, left, top, width, height, categories, series, chart_type=None):
    """Add a native (editable, right-click -> Edit Data) bar/column chart.
    `series` is a list of (name, values) tuples -- one native chart series
    per tuple. Defaults to a horizontal clustered bar (XL_CHART_TYPE.
    BAR_CLUSTERED); pass chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED for a
    vertical column chart. Returns the Chart object (not the GraphicFrame)."""
    if chart_type is None:
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    for name, values in series:
        cd.add_series(name, list(values))
    graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, cd)
    return graphic_frame.chart


def set_point_color(series, idx, colour):
    """Fill ONE data point (bar) of a chart series a solid colour, via
    python-pptx's native per-point Point.format.fill API -- e.g. tinting a
    single 'Overall' bar a distinct colour, or colouring bars by opponent
    segment. `colour` may be an RGBColor or a hex string."""
    pt = series.points[idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = _to_rgb(colour)


def colour_points_by_sign(series, values, pos_colour, neg_colour):
    """Colour every point in `series` pos_colour/neg_colour by the sign of
    the corresponding entry in `values` (>=0 -> pos_colour) -- the native
    equivalent of 11_plot_charts.py's colour_bars_by_sign, generalised to any
    series/value list rather than one specific df column."""
    for idx, v in enumerate(values):
        if v is None or v != v:      # None / NaN = a reserved blank slot, no bar to colour
            continue
        set_point_color(series, idx, pos_colour if v >= 0 else neg_colour)


def fill_series(series, colour):
    """Solid-fill an entire chart series (every point) one colour -- for a
    series that doesn't need per-point tinting (e.g. a flat grey 'breakeven'
    reference series)."""
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _to_rgb(colour)


def enable_data_labels(series, number_format='0.0"%"', font_pt=9, bold=False):
    """Turn on value-only data labels for a chart series with a given Excel
    number-format string (e.g. '0.0"%"' for a percentage, '+0.00;-0.00' for a
    signed 2dp EV figure)."""
    dl = series.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.show_value = True
    dl.show_category_name = False
    dl.show_series_name = False
    dl.show_legend_key = False
    dl.show_percentage = False
    dl.font.size = Pt(font_pt)
    dl.font.bold = bold


def set_point_label_style(series, idx, colour=None, bold=None, number_format=None, font_pt=None):
    """Style ONE data point's label, leaving the rest of the series alone.

    Used to mark a cell whose sample is below the display threshold: the bar is
    still drawn (with its wide error bar) but the number is called out so the reader
    can see at a glance that it is thin.

    PASS number_format AND font_pt whenever the series has them. python-pptx creates
    the per-point label lazily, and the label it creates does NOT inherit the
    series-level number format or font size -- so styling a point without restoring
    them silently reverts that one label to General, turning '35.0%' into '35' and
    '+2.65' into '2.654176119'. That shipped once."""
    dl = series.points[idx].data_label
    # Font FIRST, number format LAST: materialising the font rewrites the <c:dLbl>
    # child order and drops an already-written <c:numFmt>, which left the labels on
    # General format (verified in the emitted XML -- formatCode was absent).
    if font_pt is not None:
        dl.font.size = Pt(font_pt)
    if colour is not None:
        dl.font.color.rgb = _to_rgb(colour)
    if bold is not None:
        dl.font.bold = bold
    if number_format is not None:
        _force_label_number_format(dl, number_format)


def _force_label_number_format(data_label, number_format):
    """Write <c:numFmt> into a per-point <c:dLbl> by hand.

    python-pptx's DataLabel.number_format setter does not reach a per-point label --
    verified in the emitted XML, formatCode was absent every time regardless of the
    order it was set in. Same situation as add_custom_error_bars in
    14_line_optimisation.py: no working API, so raw XML.

    ELEMENT ORDER IS LOAD-BEARING. CT_DLbl's sequence is
    idx, (delete | layout, tx, numFmt, spPr, txPr, dLblPos, showLegendKey, showVal, ...)
    so numFmt must sit after idx/layout/tx and BEFORE spPr/txPr. Out of order,
    PowerPoint shows a repair prompt."""
    # NOT data_label._element -- that is the parent <c:ser>, so inserting into it
    # silently does nothing (the numFmt never reached the label).
    dLbl = data_label._get_or_add_dLbl()
    existing = dLbl.find(qn('c:numFmt'))
    if existing is not None:
        existing.set('formatCode', number_format)
        existing.set('sourceLinked', '0')
        return
    # Excel format codes carry literal double quotes ('0.0"%"'), so the attribute
    # value must escape them too -- saxutils.escape does not do quotes by default.
    code = escape(number_format, {'"': '&quot;'})
    numFmt = parse_xml(f'<c:numFmt {nsdecls("c")} formatCode="{code}" sourceLinked="0"/>')
    after = [qn('c:idx'), qn('c:layout'), qn('c:tx')]
    insert_at = 0
    for i, child in enumerate(dLbl):
        if child.tag in after:
            insert_at = i + 1
    dLbl.insert(insert_at, numFmt)


def remove_gridlines(chart):
    """Strip major/minor gridlines from both axes -- ported unchanged (same
    raw-XML technique, since python-pptx's has_major_gridlines setter can
    raise on some axis/chart-type combinations) from 11_plot_charts.py's
    identical helper."""
    for ax_tag in (qn('c:valAx'), qn('c:catAx')):
        for ax in chart._element.iter(ax_tag):
            for gl_tag in (qn('c:majorGridlines'), qn('c:minorGridlines')):
                for gl in ax.findall(gl_tag):
                    ax.remove(gl)


def set_value_axis_title(chart, text, font_pt=10):
    """Set the value (numeric) axis's title text -- the native-chart
    equivalent of a matplotlib set_xlabel()/set_ylabel() call on the
    percentage/EV axis of a horizontal bar chart."""
    ax = chart.value_axis
    ax.has_title = True
    ax.axis_title.text_frame.text = text
    for para in ax.axis_title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_pt)


def add_chapter_slide(prs, layouts, title, subtitle=None):
    """Chapter/section-divider slide, built on the template's '2_Section
    Divider' layout. That layout has no single title placeholder -- it has 5
    body placeholders (idx 10-14) designed as a highlighted table-of-contents
    (idx 10 renders with a dark navy fill + bold white text; idx 11-14 render
    with a plain light-grey fill). For v1 we only use idx 10 as the slide's
    "title" (its highlighted styling reads naturally as a chapter title) and
    idx 11 as an optional subtitle -- the fuller multi-entry TOC-highlight
    behaviour (showing all chapters with the current one highlighted) is a
    nice-to-have deferred for later, not required for template parity."""
    layout = layouts['2_Section Divider']
    slide = prs.slides.add_slide(layout)
    for idx, text in ((10, title), (11, subtitle)):
        if text is None:
            continue
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.text = text
                break
    return slide


def fig_to_stream(fig, dpi=150):
    """Render a matplotlib figure to an in-memory PNG and close it.

    Closing here matters: a deck builder emitting dozens of figures leaks them
    all into pyplot's registry otherwise, and matplotlib starts warning once
    more than 20 are open.
    """
    from io import BytesIO
    import matplotlib.pyplot as plt
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_image_slide(prs, layouts, title, fig, dpi=150):
    """Full-content-area slide holding one matplotlib figure.

    Scales to the content height, then shrinks to CONTENT_WIDTH if that made it
    too wide, and centres horizontally either way -- so charts of any aspect
    ratio sit in the same envelope as the native-chart slides.

    Deliberately minimal: takeaway lines, footnotes and chart-data registration
    differ per deck, so those stay in the deck scripts (see
    06_metagame/build_metagame_deck.py's richer local version, which this does
    NOT replace).
    """
    slide = prs.slides.add_slide(layouts['Title and Content'])
    add_title(slide, title)

    top, img_h = CONTENT_TOP, CONTENT_HEIGHT
    pic = slide.shapes.add_picture(fig_to_stream(fig, dpi=dpi), CONTENT_LEFT, top, height=img_h)
    if pic.width > CONTENT_WIDTH:
        aspect = pic.height / pic.width
        pic.width = CONTENT_WIDTH
        pic.height = int(CONTENT_WIDTH * aspect)
        pic.top = int(top + (img_h - pic.height) / 2)
    pic.left = int(CONTENT_LEFT + (CONTENT_WIDTH - pic.width) / 2)
    return slide

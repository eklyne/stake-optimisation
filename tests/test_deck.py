"""Smoke tests for the deck build.

Deliberately structural rather than visual: that the deck builds, has the six
slides in order, and that the numbers on the slides come from the same functions
the terminal prints. Nothing here can tell you a chart is ugly.
"""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

import numpy as np

from shotopt import mix
from shotopt.config import Config, Stake

try:
    from pptx import Presentation

    from shotopt import deck

    AVAILABLE = True
except ImportError:  # python-pptx / matplotlib are optional for the text commands
    AVAILABLE = False


def _config(**overrides) -> Config:
    defaults = dict(
        bankroll_eur=5000.0,
        tables=6,
        ruin_tolerance=0.01,
        rakeback_pct=0.30,
        timescale_hands=200_000,
        sim_paths=800,
        stakes=(
            Stake("50NL", 0.5, 8.18, 92.0, hands=271_592, rake_bb100=10.195,
                  current_hands=23_082),
            Stake("100NL", 1.0, 7.46, 92.0, hands=79_329, rake_bb100=8.329,
                  current_hands=11_052),
            Stake("200NL", 2.0, 4.32, 92.0, hands=38_970, rake_bb100=6.545,
                  current_hands=13_760),
            Stake("600NL", 6.0, 67.73, 92.0, hands=1_043, rake_bb100=3.857,
                  max_tables=0, current_hands=600),
        ),
    )
    defaults.update(overrides)
    return Config(**defaults)


@unittest.skipUnless(AVAILABLE, "python-pptx not installed")
class TestDeck(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.config = _config()
        cls.dir = Path(tempfile.mkdtemp())
        cls.path = deck.build(cls.config, cls.dir)
        cls.prs = Presentation(cls.path)

    def test_it_builds(self):
        self.assertTrue(self.path.is_file())
        self.assertGreater(self.path.stat().st_size, 10_000)

    def _slide(self, fragment):
        """Find a slide by a fragment of its title - indices move as the deck
        grows, and a test that breaks on reordering tests the wrong thing."""
        for slide in self.prs.slides:
            title = slide.shapes.title
            if title is not None and fragment in title.text:
                return slide
        raise AssertionError(f"no slide titled like {fragment!r}")

    def test_content_slides_in_order(self):
        titles = [
            s.shapes.title.text
            for s in self.prs.slides
            if s.shapes.title is not None and s.shapes.title.text
        ]
        expected = ["How this deck was produced",
                    "priced on its own", "know each win rate",
                    "big blinds", "euros", "one stake and nothing else",
                    "split the tables", "nearest alternatives", "optimal mix, simulated",
                    "optimal mix, twenty single lifetimes",
                    "two ways up", "bankroll grows",
                    "actually played, simulated",
                    "actually played, twenty single lifetimes",
                    "earn MORE and risk LESS", "Appendix: method"]
        self.assertEqual(len(titles), len(expected))
        for want, got in zip(expected, titles):
            self.assertIn(want, got)

    def test_every_simulation_chart_shares_one_bankroll_axis(self):
        # Autoscaled separately, the optimal and current mixes are drawn at
        # different heights and look far more alike than they are. These slides
        # exist to be compared by eye, so the frame has to be the same on all
        # four - and wide enough that nothing drawn on any of them clips.
        from shotopt import charts, sim

        best = mix.best_allocation(mix.all_allocations(self.config))
        current = mix.current_allocation(self.config)
        plotted = [a for a in (best, current) if a is not None]
        self.assertEqual(len(plotted), 2)
        results = [
            sim.simulate(self.config, a, hands=self.config.timescale_hands,
                         paths=self.config.sim_paths)
            for a in plotted
        ]
        (low, high), drawdown_xmax = charts.simulation_scales(results, self.config, plotted)

        self.assertLess(low, 0)  # room under the axis for the "broke" label
        for result, allocation in zip(results, plotted):
            self.assertLessEqual(result.checkpoint_bankroll[charts._random_rows(result)].max(),
                                 high)
            ev_end = (self.config.bankroll_eur
                      + allocation.mean_eur_per_100 * result.hands / 100)
            self.assertLessEqual(ev_end, high)
            self.assertLessEqual(float(np.percentile(result.max_drawdown, 99.5)),
                                 drawdown_xmax)

    def test_no_table_runs_off_the_slide(self):
        # The callers centre their tables, so an over-wide one loses half its
        # overflow off EACH edge - invisible in the code, obvious on the slide,
        # and exactly how the comparison table shipped 2in too wide.
        from shotopt import pptx_common as pc

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                title = slide.shapes.title
                where = title.text if title is not None else "(divider)"
                self.assertGreaterEqual(shape.left, 0, where)
                self.assertLessEqual(
                    shape.left + sum(c.width for c in shape.table.columns),
                    pc.SLIDE_WIDTH,
                    where,
                )

    def test_the_tables_price_the_horizon_in_money(self):
        # EUR/hr is the ranking number; a year's money is the one people weigh.
        header = f"{deck.timescale_label(self.config)} hands"
        best = mix.best_allocation(mix.all_allocations(self.config))
        expected = f"{deck.horizon_ev(self.config, best):,.0f}"

        for fragment in ("nearest alternatives", "bankroll grows"):
            table = next(s.table for s in self._slide(fragment).shapes if s.has_table)
            labels = [c.text for c in table.rows[0].cells]
            column = next(i for i, t in enumerate(labels) if header in t)
            column_values = [table.cell(r, column).text for r in range(1, len(table.rows))]
            self.assertIn(expected, column_values, fragment)

    def test_the_step_up_table_also_opens_with_the_played_mix(self):
        # The shared config caps its top stake, so its step-up slide is the
        # "already at the top" text version and carries no table. Build one that
        # can actually take a shot, so the comparison table is covered too.
        from shotopt import pptx_common as pc

        # A plausible ladder with room above the optimum - the shared config's
        # 600NL win rate is a deliberate outlier and would just top out.
        config = _config(stakes=(
            Stake("50NL", 0.5, 8.18, 92.0, hands=271_592, rake_bb100=10.195,
                  current_hands=23_082),
            Stake("100NL", 1.0, 7.46, 92.0, hands=79_329, rake_bb100=8.329,
                  current_hands=11_052),
            Stake("200NL", 2.0, 4.32, 92.0, hands=38_970, rake_bb100=6.545,
                  current_hands=13_760),
            Stake("400NL", 4.0, 4.41, 92.0, hands=11_931, rake_bb100=4.411,
                  current_hands=4_632),
        ))
        self.assertTrue(mix.step_up_options(
            config, mix.best_allocation(mix.all_allocations(config))))

        prs = Presentation(deck.build(config, Path(tempfile.mkdtemp())))
        slide = next(s for s in prs.slides if s.shapes.title is not None
                     and "two ways up" in s.shapes.title.text)
        table = next(s.table for s in slide.shapes if s.has_table)

        def fill(row):
            found = table.cell(row, 0)._tc.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
            )
            return None if found is None else found.get("val")

        self.assertIn(mix.current_allocation(config).label,
                      [c.text for c in table.rows[1].cells])
        self.assertEqual(fill(1), str(pc.COL_ORANGE))   # played mix
        self.assertEqual(fill(2), str(pc.COL_GREEN))    # the optimum

    def test_four_section_dividers(self):
        # Chapter slides use a layout with no title placeholder, so they show up
        # as the untitled slides. Four sections, four dividers.
        untitled = [
            s for s in self.prs.slides
            if s.shapes.title is None or not s.shapes.title.text
        ]
        self.assertEqual(len(untitled), 4)

    def test_the_run_parameters_slide_reports_the_config_it_ran_on(self):
        slide = self._slide("How this deck was produced")
        settings, stakes = (s.table for s in slide.shapes if s.has_table)

        values = {settings.cell(r, 0).text: settings.cell(r, 1).text
                  for r in range(1, len(settings.rows))}
        self.assertIn("5,000", values["Bankroll"])
        self.assertIn(str(self.config.tables), values["Tables played at once"])
        self.assertIn(f"{self.config.timescale_hands:,}", values["Simulation timescale"])
        self.assertIn(f"{self.config.sim_paths:,}", values["Lifetimes simulated"])

        # Every stake, including the one excluded by max_tables = 0 - the slide
        # documents the inputs, not the ones that survived screening.
        names = [stakes.cell(r, 0).text for r in range(1, len(stakes.rows))]
        self.assertEqual(names, [s.name for s in self.config.stakes])

    def test_the_run_parameters_slide_follows_a_cli_override(self):
        # The value that ran, not the value in the file - a --bankroll override
        # must move the slide or the deck can contradict its own cover.
        config = _config(bankroll_eur=20_000.0)
        prs = Presentation(deck.build(config, Path(tempfile.mkdtemp())))
        slide = next(s for s in prs.slides if s.shapes.title is not None
                     and "How this deck was produced" in s.shapes.title.text)
        settings = next(s.table for s in slide.shapes if s.has_table)
        values = {settings.cell(r, 0).text: settings.cell(r, 1).text
                  for r in range(1, len(settings.rows))}
        self.assertIn("20,000", values["Bankroll"])

    def test_the_stake_table_lists_every_stake_including_excluded(self):
        table = next(s.table for s in self._slide("priced on its own").shapes if s.has_table)
        names = [table.cell(r, 0).text for r in range(1, len(table.rows))]
        self.assertEqual(names, [s.name for s in self.config.stakes])

    def test_the_excluded_stake_is_marked_as_such_on_the_slide(self):
        table = next(s.table for s in self._slide("priced on its own").shapes if s.has_table)
        verdicts = {
            table.cell(r, 0).text: table.cell(r, len(table.columns) - 1).text
            for r in range(1, len(table.rows))
        }
        self.assertEqual(verdicts["50NL"], "in the mix")
        self.assertIn("max_tables", verdicts["600NL"])

    def test_the_configuration_table_centres_on_the_chosen_mix(self):
        table = next(
            s.table for s in self._slide("nearest alternatives").shapes if s.has_table
        )
        last = len(table.columns) - 1  # the note column, wherever it has ended up
        notes = [table.cell(r, last).text for r in range(1, len(table.rows))]
        self.assertEqual(sum("CHOSEN" in n for n in notes), 1)

        best = mix.best_allocation(mix.all_allocations(self.config))
        labels = [table.cell(r, 0).text for r in range(1, len(table.rows))]
        self.assertIn(best.label, labels)

    def test_every_table_opens_with_the_played_mix_as_its_benchmark(self):
        # The played mix has no slide of its own any more - it is the first data
        # row of every allocation table, so each one can be read against what is
        # actually being done today.
        from shotopt import pptx_common as pc

        current = mix.current_allocation(self.config)
        for fragment in ("one stake and nothing else", "nearest alternatives",
                         "bankroll grows"):
            table = next(s.table for s in self._slide(fragment).shapes if s.has_table)
            row = [c.text for c in table.rows[1].cells]
            self.assertIn(current.label, row, fragment)
            fill = table.cell(1, 0)._tc.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
            )
            self.assertEqual(fill.get("val"), str(pc.COL_ORANGE), fragment)

    def test_the_optimum_is_highlighted_green_where_it_appears(self):
        from shotopt import pptx_common as pc

        best = mix.best_allocation(mix.all_allocations(self.config))
        for fragment in ("nearest alternatives",):
            table = next(s.table for s in self._slide(fragment).shapes if s.has_table)
            green = [
                r for r in range(1, len(table.rows))
                if (lambda f: f is not None and f.get("val") == str(pc.COL_GREEN))(
                    table.cell(r, 0)._tc.find(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                    )
                )
            ]
            self.assertEqual(len(green), 1, fragment)
            self.assertIn(best.label, [c.text for c in table.rows[green[0]].cells], fragment)

    def test_the_waterfall_reconciles_to_the_screen(self):
        # before rake - rake + rakeback must equal what the screen calls banked,
        # or the slide is telling a different story from the terminal.
        for screen in mix.screen_stakes(self.config):
            if not screen.kept:
                continue
            stake = screen.stake
            before = stake.winrate_bb100 + stake.rake_bb100
            banked = before - stake.rake_bb100 + stake.rake_bb100 * self.config.rakeback_pct
            self.assertAlmostEqual(banked, screen.mean_eur_per_100 / stake.bb_eur, places=9)

    def test_it_builds_with_no_stake_inside_tolerance(self):
        # No chosen mix means nothing to step up FROM, nothing to simulate and
        # nothing to compare the current split against, so sections 2 and 3 drop
        # out rather than inventing a subject.
        config = _config(bankroll_eur=100.0, ruin_tolerance=1e-9)
        self.assertIsNone(mix.best_allocation(mix.all_allocations(config)))
        prs = Presentation(deck.build(config, Path(tempfile.mkdtemp())))
        titles = [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]
        self.assertFalse(any("actually playing" in t for t in titles))

    def test_it_builds_without_a_current_distribution(self):
        # current_hands is optional; section 3 simply does not appear.
        stakes = tuple(
            Stake(s.name, s.bb_eur, s.winrate_bb100, s.stdev_bb100, s.hands,
                  s.max_tables, s.rake_bb100)
            for s in _config().stakes
        )
        config = _config(stakes=stakes)
        self.assertIsNone(mix.current_allocation(config))
        prs = Presentation(deck.build(config, Path(tempfile.mkdtemp())))
        titles = [
            s.shapes.title.text for s in prs.slides
            if s.shapes.title is not None and s.shapes.title.text
        ]
        self.assertFalse(any("actually play" in t for t in titles))


if __name__ == "__main__":
    unittest.main()
